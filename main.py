"""
Nepal Energy Agentic AI Autonomous Demo - Flat, beginner-friendly version.

One file contains synthetic CSV generation, multi-agent workflow, lightweight
Agentic RAG, evaluation, charting, PowerPoint, and DOCX report generation.

Run:
    python main.py

Then launch Streamlit:
    streamlit run streamlit_app.py
"""
from __future__ import annotations

import json
import math
import os
import random
from dataclasses import dataclass
from datetime import datetime
from pathlib import Path
from typing import Dict, List, Tuple

import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Inches, Pt, RGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PptRGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PptInches, Pt as PptPt
from sklearn.ensemble import IsolationForest
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity

ROOT = Path(__file__).resolve().parent
TODAY = datetime.now().strftime("%Y-%m-%d")

# Flat filenames: no project subfolders required.
CUSTOMERS_CSV = ROOT / "synthetic_customers.csv"
FEEDERS_CSV = ROOT / "synthetic_feeders.csv"
HYDRO_CSV = ROOT / "synthetic_hydro_trade.csv"
WORKFORCE_CSV = ROOT / "synthetic_workforce_tasks.csv"
KNOWLEDGE_CSV = ROOT / "synthetic_energy_knowledge.csv"
RISK_CSV = ROOT / "risk_scores.csv"
FEEDER_PRIORITY_CSV = ROOT / "feeder_priority_scores.csv"
RAG_EVAL_CSV = ROOT / "rag_evaluation.csv"
BUSINESS_IMPACT_CSV = ROOT / "business_impact_summary.csv"
AGENT_LOG_JSON = ROOT / "agent_chat_log.json"
PPTX_FILE = ROOT / "nepal_energy_agentic_ai_demo.pptx"
DOCX_FILE = ROOT / "nepal_energy_agentic_ai_report.docx"
README_OUTPUT = ROOT / "RUNNING_STEPS.md"

CHART_RISK = ROOT / "chart_risk_distribution.png"
CHART_REVENUE = ROOT / "chart_revenue_opportunity.png"
CHART_FEEDER = ROOT / "chart_feeder_priority.png"
CHART_WORKFORCE = ROOT / "chart_workforce_redeployment.png"
CHART_RAG = ROOT / "chart_rag_metrics.png"

PUBLIC_CONTEXT = [
    {
        "name": "NEA Annual Report 2024/2025",
        "url": "https://nea.org.np/detail/annual-report-20242025",
        "note": "NEA publishes annual energy, infrastructure, operations, and planning updates.",
    },
    {
        "name": "Reuters: Nepal electricity export to Bangladesh",
        "url": "https://www.reuters.com/business/energy/nepal-begins-first-power-exports-bangladesh-via-indias-grid-2025-06-16/",
        "note": "Nepal began electricity exports to Bangladesh via India's grid, showing cross-border trade opportunity.",
    },
    {
        "name": "GitHub electricity theft / anomaly examples",
        "url": "https://github.com/henryRDlab/ElectricityTheftDetection",
        "note": "Open-source electricity theft and anomaly detection examples show related technical feasibility.",
    },
    {
        "name": "World Bank Nepal energy infrastructure assessment",
        "url": "https://documents1.worldbank.org/curated/en/592481554093658883/pdf/Nepal-Energy-Infrastructure-Sector-Assessment.pdf",
        "note": "Energy infrastructure reports support the need for data-driven planning and grid modernization.",
    },
]

DISTRICTS = ["Kathmandu", "Lalitpur", "Bhaktapur", "Pokhara", "Chitwan", "Biratnagar", "Butwal", "Nepalgunj", "Dharan", "Janakpur", "Hetauda", "Birgunj"]
TARIFFS = ["Domestic", "Commercial", "Industrial", "Irrigation", "EV Charging"]
FEEDER_NAMES = [f"FD-{i:03d}" for i in range(1, 31)]


def money(n: float) -> str:
    return f"NPR {n:,.0f}"


def pct(n: float) -> str:
    return f"{n:.1f}%"


def log_event(logs: List[Dict], agent: str, message: str, data: Dict | None = None):
    logs.append({
        "time": datetime.now().isoformat(timespec="seconds"),
        "agent": agent,
        "message": message,
        "data": data or {},
    })


def create_synthetic_csvs(seed: int = 42) -> None:
    """Create multiple synthetic Nepal energy CSV files in the root folder."""
    random.seed(seed)
    np.random.seed(seed)

    customers = []
    for i in range(1, 801):
        district = random.choice(DISTRICTS)
        tariff = np.random.choice(TARIFFS, p=[0.55, 0.2, 0.12, 0.08, 0.05])
        feeder = random.choice(FEEDER_NAMES)
        base = {
            "Domestic": np.random.normal(120, 35),
            "Commercial": np.random.normal(550, 180),
            "Industrial": np.random.normal(3500, 1300),
            "Irrigation": np.random.normal(900, 300),
            "EV Charging": np.random.normal(2200, 700),
        }[tariff]
        prev_kwh = max(20, base)
        # Inject suspicious drops/spikes for a subset.
        suspicious = np.random.rand() < 0.12
        if suspicious:
            curr_kwh = prev_kwh * np.random.uniform(0.05, 0.45) if np.random.rand() < 0.75 else prev_kwh * np.random.uniform(2.0, 3.8)
            tamper = np.random.choice(["seal_broken", "reverse_flow", "meter_offline", "bypass_suspected"])
        else:
            curr_kwh = prev_kwh * np.random.uniform(0.82, 1.18)
            tamper = np.random.choice(["normal", "normal", "normal", "late_reading"])
        tariff_rate = {"Domestic": 10.5, "Commercial": 13.0, "Industrial": 11.5, "Irrigation": 8.0, "EV Charging": 12.0}[tariff]
        outstanding = max(0, np.random.gamma(2, 1800) - (0 if suspicious else 1200))
        meter_age = np.random.randint(1, 16)
        outage_complaints = np.random.poisson(1.0 if not suspicious else 2.8)
        location_lat = 26.4 + random.random() * 3.7
        location_lon = 80.0 + random.random() * 8.5
        customers.append({
            "customer_id": f"CUST-{i:05d}",
            "district": district,
            "feeder_id": feeder,
            "tariff_category": tariff,
            "previous_month_kwh": round(prev_kwh, 2),
            "current_month_kwh": round(curr_kwh, 2),
            "kwh_change_pct": round((curr_kwh - prev_kwh) / max(prev_kwh, 1) * 100, 2),
            "tariff_rate_npr_per_kwh": tariff_rate,
            "monthly_bill_npr": round(curr_kwh * tariff_rate, 2),
            "outstanding_balance_npr": round(outstanding, 2),
            "meter_age_years": meter_age,
            "meter_event": tamper,
            "outage_complaints_90d": outage_complaints,
            "latitude": round(location_lat, 5),
            "longitude": round(location_lon, 5),
            "synthetic_label_suspicious": int(suspicious),
        })
    pd.DataFrame(customers).to_csv(CUSTOMERS_CSV, index=False)

    feeders = []
    for feeder in FEEDER_NAMES:
        district = random.choice(DISTRICTS)
        input_mwh = np.random.uniform(700, 5000)
        loss_pct = np.random.uniform(6, 34)
        sold_mwh = input_mwh * (1 - loss_pct / 100)
        outages = np.random.poisson(4 + max(0, loss_pct - 15) / 4)
        overloaded = np.random.uniform(0.55, 1.35)
        feeders.append({
            "feeder_id": feeder,
            "district": district,
            "energy_input_mwh": round(input_mwh, 2),
            "energy_sold_mwh": round(sold_mwh, 2),
            "technical_nontechnical_loss_pct": round(loss_pct, 2),
            "outage_events_30d": int(outages),
            "peak_loading_ratio": round(overloaded, 2),
            "field_visits_available_month": np.random.randint(8, 25),
        })
    pd.DataFrame(feeders).to_csv(FEEDERS_CSV, index=False)

    months = pd.date_range("2025-01-01", periods=18, freq="MS")
    hydro = []
    for m in months:
        wet = 1 if m.month in [6, 7, 8, 9, 10] else 0
        generation = np.random.normal(2300 if wet else 1200, 220)
        demand = np.random.normal(1500 if wet else 1750, 160)
        export_mwh = max(0, generation - demand) * np.random.uniform(0.25, 0.55)
        import_mwh = max(0, demand - generation) * np.random.uniform(0.35, 0.65)
        export_price = np.random.uniform(5.2, 8.8)
        hydro.append({
            "month": m.strftime("%Y-%m"),
            "hydro_generation_gwh": round(max(700, generation), 2),
            "domestic_demand_gwh": round(max(900, demand), 2),
            "export_energy_gwh": round(export_mwh, 2),
            "import_energy_gwh": round(import_mwh, 2),
            "estimated_export_price_npr_per_kwh": round(export_price, 2),
            "ev_growth_index": round(100 + (m.to_period('M').ordinal - months[0].to_period('M').ordinal) * 4.2 + np.random.normal(0, 3), 2),
        })
    pd.DataFrame(hydro).to_csv(HYDRO_CSV, index=False)

    # Volumes are synthetic but scaled like an enterprise pilot across several distribution centers.
    # They create a credible productivity demo: dozens of FTE-equivalent hours can be redeployed
    # from repetitive review into field execution, governance, and customer service.
    workforce = pd.DataFrame([
        ["meter_reading_validation", "Manual meter read review", 28000, 12, 0.72],
        ["billing_exception_review", "Billing exception investigation", 15000, 20, 0.68],
        ["field_visit_prioritization", "Field inspection planning", 3500, 35, 0.60],
        ["customer_complaint_triage", "Complaint triage and routing", 22000, 14, 0.58],
        ["report_generation", "Monthly stakeholder reporting", 400, 180, 0.82],
        ["document_search", "Policy and evidence search", 11000, 10, 0.74],
        ["loss_analysis", "Feeder loss analysis", 2500, 45, 0.66],
    ], columns=["task_id", "task_name", "monthly_volume", "minutes_per_case_manual", "automation_assist_rate"])
    workforce.to_csv(WORKFORCE_CSV, index=False)

    knowledge = pd.DataFrame([
        ["K1", "NEA Annual Reporting", "NEA annual reports describe system expansion, generation, transmission, distribution, customer service, and financial performance. These reports are useful evidence for stakeholder AI planning.", "NEA Annual Report 2024/2025", "https://nea.org.np/detail/annual-report-20242025"],
        ["K2", "Cross-border Export", "Nepal's growing hydropower base creates a business case for export timing, wet-season surplus management, and cross-border trade planning with India and Bangladesh.", "Reuters electricity export to Bangladesh", "https://www.reuters.com/business/energy/nepal-begins-first-power-exports-bangladesh-via-indias-grid-2025-06-16/"],
        ["K3", "Loss Reduction", "Distribution utilities can reduce non-technical losses by ranking high-risk customers, comparing expected and observed usage, prioritizing field visits, and tracking recovery after inspection.", "Synthetic utility operations note", "local synthetic knowledge"],
        ["K4", "Agentic RAG", "Agentic RAG improves over basic chatbot search by letting specialized agents retrieve evidence, reason over tables, check confidence, calculate evaluation metrics, and generate reports autonomously.", "Architecture note", "local synthetic knowledge"],
        ["K5", "Workforce Productivity", "Automation should be framed as workforce capacity redeployment. AI can reduce repetitive manual workload while keeping humans for approvals, safety, field execution, and governance.", "Change management note", "local synthetic knowledge"],
        ["K6", "Evaluation", "RAG quality can be measured with MRR, retrieval relevance, semantic similarity, groundedness, faithfulness, and hallucination risk based on unsupported claims.", "Evaluation note", "local synthetic knowledge"],
        ["K7", "Open-source feasibility", "Open-source projects demonstrate electricity theft detection, time-series anomaly detection, and Agentic RAG patterns that can be adapted for a Nepal energy proof of concept.", "GitHub examples", "https://github.com/henryRDlab/ElectricityTheftDetection"],
    ], columns=["doc_id", "topic", "text", "source_name", "source_url"])
    knowledge.to_csv(KNOWLEDGE_CSV, index=False)


def build_contextual_chunks(knowledge: pd.DataFrame) -> pd.DataFrame:
    rows = []
    for _, r in knowledge.iterrows():
        text = str(r["text"])
        topic = str(r["topic"])
        source = str(r["source_name"])
        words = text.split()
        # Very simple chunking for a demo; contextual prefix improves retrieval.
        for i in range(0, len(words), 42):
            chunk = " ".join(words[i:i+70])
            if not chunk:
                continue
            contextual = f"Topic: {topic}. Source: {source}. Nepal energy context. {chunk}"
            rows.append({
                "chunk_id": f"{r['doc_id']}-CH{i//42+1}",
                "doc_id": r["doc_id"],
                "topic": topic,
                "source_name": source,
                "source_url": r["source_url"],
                "chunk_text": chunk,
                "contextual_text": contextual,
            })
    return pd.DataFrame(rows)


@dataclass
class RagIndex:
    chunks: pd.DataFrame
    vectorizer: TfidfVectorizer
    matrix: object


def build_rag_index(knowledge: pd.DataFrame) -> RagIndex:
    chunks = build_contextual_chunks(knowledge)
    vectorizer = TfidfVectorizer(ngram_range=(1, 2), stop_words="english", min_df=1)
    matrix = vectorizer.fit_transform(chunks["contextual_text"].tolist())
    return RagIndex(chunks=chunks, vectorizer=vectorizer, matrix=matrix)


def retrieve(index: RagIndex, query: str, top_k: int = 4) -> pd.DataFrame:
    qv = index.vectorizer.transform([query])
    sims = cosine_similarity(qv, index.matrix)[0]
    order = np.argsort(-sims)[:top_k]
    out = index.chunks.iloc[order].copy()
    out["similarity"] = sims[order]
    return out


def generate_grounded_answer(query: str, retrieved: pd.DataFrame) -> str:
    bullets = []
    for _, row in retrieved.head(3).iterrows():
        bullets.append(f"- {row['topic']}: {row['chunk_text']} [Source: {row['source_name']}]")
    return (
        f"Question: {query}\n\n"
        "Grounded answer from retrieved evidence:\n"
        + "\n".join(bullets)
        + "\n\nRecommended executive action: launch a controlled pilot that combines anomaly scoring, feeder prioritization, field visit optimization, and monthly evidence-based reporting."
    )


def evaluate_rag(index: RagIndex) -> pd.DataFrame:
    eval_questions = [
        ("How can NEA reduce non-technical losses with AI?", "K3"),
        ("Why is cross-border electricity trade important for Nepal?", "K2"),
        ("How should workforce reduction be framed for stakeholders?", "K5"),
        ("Which metrics evaluate RAG quality?", "K6"),
        ("Why is Agentic RAG better than a chatbot?", "K4"),
    ]
    rows = []
    for q, gold_doc in eval_questions:
        top = retrieve(index, q, top_k=5)
        doc_ids = top["doc_id"].tolist()
        rank = doc_ids.index(gold_doc) + 1 if gold_doc in doc_ids else None
        mrr = 1.0 / rank if rank else 0.0
        avg_similarity = float(top["similarity"].mean())
        relevance = float(top.iloc[0]["similarity"])
        answer = generate_grounded_answer(q, top)
        answer_terms = set(w.lower().strip(".,:;[]()") for w in answer.split() if len(w) > 5)
        context_terms = set(w.lower().strip(".,:;[]()") for w in " ".join(top["chunk_text"].tolist()).split() if len(w) > 5)
        supported = len(answer_terms.intersection(context_terms))
        total = max(1, len(answer_terms))
        groundedness = supported / total
        hallucination_rate = max(0.02, 1 - groundedness)  # simple demo estimate
        rows.append({
            "question": q,
            "gold_doc_id": gold_doc,
            "rank_of_gold_doc": rank if rank else "not_found",
            "mrr": round(mrr, 3),
            "relevance_top1": round(relevance, 3),
            "avg_similarity_top5": round(avg_similarity, 3),
            "groundedness_estimate": round(groundedness, 3),
            "hallucination_rate_estimate": round(hallucination_rate, 3),
            "answer_preview": answer[:240] + "...",
        })
    return pd.DataFrame(rows)


def run_customer_risk(customers: pd.DataFrame) -> pd.DataFrame:
    df = customers.copy()
    event_score = df["meter_event"].map({"normal": 0, "late_reading": 0.15, "seal_broken": 0.8, "reverse_flow": 1.0, "meter_offline": 0.75, "bypass_suspected": 1.0}).fillna(0.2)
    features = df[["previous_month_kwh", "current_month_kwh", "kwh_change_pct", "outstanding_balance_npr", "meter_age_years", "outage_complaints_90d"]].copy()
    features["event_score"] = event_score
    model = IsolationForest(n_estimators=100, contamination=0.13, random_state=7)
    pred = model.fit_predict(features)
    anomaly_raw = -model.score_samples(features)
    scaled = (anomaly_raw - anomaly_raw.min()) / max(1e-9, (anomaly_raw.max() - anomaly_raw.min()))
    drop_component = np.clip((-df["kwh_change_pct"]) / 100, 0, 1)
    event_component = event_score
    balance_component = np.clip(df["outstanding_balance_npr"] / 15000, 0, 1)
    risk_score = 100 * (0.42 * scaled + 0.28 * drop_component + 0.2 * event_component + 0.1 * balance_component)
    expected_kwh = df["previous_month_kwh"] * 0.96
    lost_kwh_month = np.clip(expected_kwh - df["current_month_kwh"], 0, None)
    df["ai_anomaly_flag"] = (pred == -1).astype(int)
    df["risk_score_0_100"] = np.round(risk_score, 2)
    df["estimated_lost_kwh_month"] = np.round(lost_kwh_month, 2)
    df["estimated_revenue_recovery_npr_year"] = np.round(lost_kwh_month * df["tariff_rate_npr_per_kwh"] * 12 * (df["risk_score_0_100"] / 100), 2)
    df["recommended_action"] = np.where(df["risk_score_0_100"] >= 70, "Immediate inspection", np.where(df["risk_score_0_100"] >= 45, "Remote validation / call", "Monitor"))
    return df.sort_values("risk_score_0_100", ascending=False)


def run_feeder_priority(feeders: pd.DataFrame, risk: pd.DataFrame) -> pd.DataFrame:
    agg = risk.groupby("feeder_id").agg(
        suspicious_customers=("ai_anomaly_flag", "sum"),
        revenue_opportunity_npr=("estimated_revenue_recovery_npr_year", "sum"),
        avg_customer_risk=("risk_score_0_100", "mean"),
    ).reset_index()
    df = feeders.merge(agg, on="feeder_id", how="left").fillna(0)
    df["priority_score_0_100"] = np.round(
        0.35 * df["technical_nontechnical_loss_pct"] / df["technical_nontechnical_loss_pct"].max() * 100
        + 0.25 * df["suspicious_customers"] / max(1, df["suspicious_customers"].max()) * 100
        + 0.25 * df["revenue_opportunity_npr"] / max(1, df["revenue_opportunity_npr"].max()) * 100
        + 0.15 * np.clip(df["peak_loading_ratio"], 0, 1.4) / 1.4 * 100,
        2,
    )
    df["field_visit_recommendation"] = np.where(df["priority_score_0_100"] >= 70, "Deploy team this week", np.where(df["priority_score_0_100"] >= 50, "Schedule this month", "Monitor digitally"))
    return df.sort_values("priority_score_0_100", ascending=False)


def calculate_business_impact(risk: pd.DataFrame, workforce: pd.DataFrame, hydro: pd.DataFrame) -> pd.DataFrame:
    annual_revenue_recovery = float(risk.query("recommended_action == 'Immediate inspection'")["estimated_revenue_recovery_npr_year"].sum())
    monthly_revenue_recovery = annual_revenue_recovery / 12
    scaled_25000_customer_recovery = annual_revenue_recovery * (25000 / max(1, len(risk)))

    wf = workforce.copy()
    wf["manual_hours_month"] = wf["monthly_volume"] * wf["minutes_per_case_manual"] / 60
    wf["ai_assisted_hours_month"] = wf["manual_hours_month"] * (1 - wf["automation_assist_rate"])
    baseline_hours = float(wf["manual_hours_month"].sum())
    assisted_hours = float(wf["ai_assisted_hours_month"].sum())
    saved_hours = baseline_hours - assisted_hours
    baseline_fte = baseline_hours / 160
    remaining_fte = assisted_hours / 160
    redeployable_fte = saved_hours / 160
    half_workload_target = baseline_fte * 0.5

    export_revenue = float((hydro["export_energy_gwh"] * 1_000_000 * hydro["estimated_export_price_npr_per_kwh"]).sum())

    rows = [
        ["Annual revenue recovery opportunity from high-risk customers", annual_revenue_recovery, "NPR/year", "Estimated from synthetic lost kWh, tariff, and AI risk score."],
        ["Monthly revenue recovery opportunity", monthly_revenue_recovery, "NPR/month", "Annual opportunity divided by 12."],
        ["Scaled annual recovery for 25,000-customer pilot", scaled_25000_customer_recovery, "NPR/year", "Linear scale-up from this 800-customer synthetic portfolio; for demo planning only."],
        ["Total modeled export revenue opportunity", export_revenue, "NPR over synthetic period", "Based on synthetic hydro/export CSV, not an official forecast."],
        ["Baseline repetitive manual work", baseline_fte, "FTE-equivalent/month", "Manual review, reporting, complaint triage, loss analysis, field planning."],
        ["Conservative half-workload redeployment target", half_workload_target, "FTE-equivalent/month", "Stakeholder framing: capacity redeployment, not layoffs."],
        ["Modeled maximum redeployable capacity", redeployable_fte, "FTE-equivalent/month", "AI-assisted hours saved divided by 160 hours per FTE-month."],
        ["Remaining manual work after AI assistance", remaining_fte, "FTE-equivalent/month", "Humans remain needed for approvals, field work, governance, and safety."],
        ["Modeled repetitive workload reduction", saved_hours / baseline_hours * 100, "%", "Estimated productivity improvement across repetitive tasks."],
    ]
    return pd.DataFrame(rows, columns=["metric", "value", "unit", "explanation"])


def make_charts(risk: pd.DataFrame, feeder: pd.DataFrame, impact: pd.DataFrame, rag_eval: pd.DataFrame) -> None:
    plt.figure(figsize=(8, 4.5))
    plt.hist(risk["risk_score_0_100"], bins=20)
    plt.title("Customer suspicious activity risk distribution")
    plt.xlabel("Risk score (0-100)")
    plt.ylabel("Customers")
    plt.tight_layout(); plt.savefig(CHART_RISK, dpi=160); plt.close()

    top_rev = risk.head(15).sort_values("estimated_revenue_recovery_npr_year")
    plt.figure(figsize=(8, 5))
    plt.barh(top_rev["customer_id"], top_rev["estimated_revenue_recovery_npr_year"])
    plt.title("Top revenue recovery opportunities")
    plt.xlabel("NPR/year")
    plt.tight_layout(); plt.savefig(CHART_REVENUE, dpi=160); plt.close()

    top_feeder = feeder.head(10).sort_values("priority_score_0_100")
    plt.figure(figsize=(8, 5))
    plt.barh(top_feeder["feeder_id"], top_feeder["priority_score_0_100"])
    plt.title("Feeder priority score")
    plt.xlabel("Priority score (0-100)")
    plt.tight_layout(); plt.savefig(CHART_FEEDER, dpi=160); plt.close()

    vals = impact.set_index("metric")["value"]
    labels = ["Baseline FTE", "Remaining FTE", "Redeployable FTE"]
    values = [vals["Baseline repetitive manual work"], vals["Remaining manual work after AI assistance"], vals["Modeled maximum redeployable capacity"]]
    plt.figure(figsize=(8, 4.5))
    plt.bar(labels, values)
    plt.title("Workforce productivity / redeployment estimate")
    plt.ylabel("FTE-equivalent/month")
    plt.tight_layout(); plt.savefig(CHART_WORKFORCE, dpi=160); plt.close()

    metrics = ["mrr", "relevance_top1", "avg_similarity_top5", "groundedness_estimate"]
    plt.figure(figsize=(8, 4.5))
    plt.bar(metrics, [rag_eval[m].mean() for m in metrics])
    plt.title("Average RAG evaluation metrics")
    plt.ylim(0, 1)
    plt.tight_layout(); plt.savefig(CHART_RAG, dpi=160); plt.close()


def add_title(slide, title, subtitle=None):
    box = slide.shapes.add_textbox(PptInches(0.45), PptInches(0.25), PptInches(12.4), PptInches(0.55))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = PptPt(24)
    p.font.bold = True
    p.font.color.rgb = PptRGBColor(20, 56, 92)
    if subtitle:
        box2 = slide.shapes.add_textbox(PptInches(0.48), PptInches(0.8), PptInches(12), PptInches(0.3))
        q = box2.text_frame.paragraphs[0]
        q.text = subtitle
        q.font.size = PptPt(11)
        q.font.color.rgb = PptRGBColor(90, 90, 90)


def add_bullets(slide, x, y, w, h, bullets, font_size=15):
    box = slide.shapes.add_textbox(PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    tf = box.text_frame
    tf.clear()
    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = 0
        p.font.size = PptPt(font_size)
        p.space_after = PptPt(8)



AGENT_BLUEPRINT = [
    ("1","Orchestrator Agent","Executive workflow coordinator","Coordinates the full autonomous run, triggers agents, tracks outputs, and ensures PPTX/DOCX/CSV/chart artifacts are created.","Executives, PMO, digital transformation team","Reduces status chasing, repeated coordination, and manual reporting follow-up.","Communicates with every agent and acts as the central controller."),
    ("2","Data Agent","Data preparation specialist","Creates/loads customer, feeder, hydropower, workforce, and knowledge CSV files and validates columns.","IT data team, billing data team, MIS staff","Reduces repetitive spreadsheet cleaning, joining, and reformatting.","Passes clean tables to Risk, Feeder, Workforce, and RAG agents."),
    ("3","Customer Risk Agent","Revenue protection analyst","Scores suspicious customer behavior using kWh drops, meter events, arrears, complaints, tariff class, and anomaly scores.","Revenue protection unit, distribution managers, field supervisors","Reduces manual scanning of thousands of accounts and prioritizes high-risk cases.","Receives clean data from Data Agent and sends ranked customers to Feeder and Business Impact agents."),
    ("4","Feeder Prioritization Agent","Distribution planning analyst","Ranks feeders using loss, overload, outage, suspicious customer density, and recovery opportunity.","Distribution managers, operations engineers, planning teams","Reduces manual feeder-by-feeder comparison and field-crew prioritization.","Combines customer risk with feeder data and sends priorities to dashboards/reports."),
    ("5","Hydropower & Trade Agent","Energy planning analyst","Adds hydropower seasonality, demand, import/export, EV-load, and power-trade context.","System planning, power trade, dispatch, generation planning","Reduces manual scenario preparation for seasonal and export/import discussion.","Provides strategic context to Business Impact and Reporting agents."),
    ("6","Workforce Productivity Agent","Process automation analyst","Estimates baseline repetitive FTE-equivalent workload, AI-assisted workload, and redeployable capacity.","HR, department heads, managers, process owners","Shows how repetitive work can be reduced while redeploying staff to higher-value work.","Uses workforce task data and impact assumptions to feed charts/reports."),
    ("7","Agentic RAG Agent","Knowledge retrieval specialist","Chunks knowledge, adds contextual metadata, creates local embeddings, retrieves evidence, and writes grounded answers.","Executives, policy teams, IT, legal, training teams","Reduces time spent searching reports, policies, and previous decision notes.","Retrieves evidence for Evaluation and Reporting agents."),
    ("8","RAG Evaluation Agent","AI quality auditor","Measures MRR, relevance, similarity, groundedness, and hallucination-risk estimates.","AI governance team, IT QA, audit, management","Reduces manual AI quality checking and increases trust through transparent metrics.","Consumes retrieved evidence/answers and reports metrics to dashboard/PPTX/DOCX."),
    ("9","Visualization Agent","Dashboard designer","Creates charts for risk, revenue, feeder priority, workforce redeployment, and RAG quality.","Managers, executives, reporting analysts","Reduces manual charting and screenshot work.","Receives final tables and creates visual assets."),
    ("10","Reporting Agent","Stakeholder communication specialist","Creates stakeholder-ready PowerPoint, DOCX report, CSV outputs, and logs in one run.","Executives, managers, external stakeholders","Reduces report writing, formatting, and repeated slide preparation.","Final receiver of outputs from all agents."),
]
TECH_STACK = [
    ("Frontend","Streamlit","Interactive public dashboard; Streamlit Community Cloud entry point is streamlit_app.py."),
    ("Language","Python","Simple flat project with main.py and streamlit_app.py."),
    ("LLM used","Local deterministic demo LLM / grounded template generator","No paid API key. It writes answers only from retrieved evidence for safer public demo behavior."),
    ("Production LLM options","Gemini, OpenAI, Groq, Ollama Llama/Mistral","Can replace the local generator later while keeping the same RAG/evaluation design."),
    ("RAG type","Agentic RAG","Retrieval is part of an autonomous workflow: retrieve, answer, evaluate, chart, and report."),
    ("Hybrid search","Contextual lexical-semantic retrieval","TF-IDF n-grams plus topic/source/Nepal-energy prefixes approximate hybrid behavior in a lightweight way."),
    ("Embedding","Contextual TF-IDF embeddings","Each chunk is embedded after adding topic, source, and Nepal energy context. Production can use text-embedding-004, BGE, E5, or OpenAI embeddings."),
    ("Vector store","Local sparse vector matrix","Runs easily on Streamlit Cloud. Production can use Chroma, FAISS, Pinecone, Weaviate, Milvus, or pgvector."),
    ("Anomaly ML","Isolation Forest + rules","Detects unusual consumption and combines it with interpretable energy-sector risk signals."),
    ("Analytics","Pandas, NumPy, scikit-learn","Creates scores, rankings, evaluation metrics, and business-impact estimates."),
    ("Visualization","Matplotlib","Generates charts for Streamlit, PPTX, and DOCX."),
    ("Reports","python-pptx and python-docx","Creates stakeholder PowerPoint and long Word report automatically."),
]
RAG_METRIC_DEFS = [
    ("MRR","Mean Reciprocal Rank","Checks whether reference evidence appears near the top of retrieval results. Higher is better."),
    ("Top-1 relevance","Similarity score of best retrieved chunk","Measures whether the first evidence is strongly related to the question."),
    ("Average similarity top-5","Mean similarity of retrieved evidence","Shows how coherent the evidence set is across top results."),
    ("Groundedness","Share of answer terms supported by evidence","Estimates whether answers are supported by retrieved chunks."),
    ("Hallucination-risk estimate","1 - groundedness with a small floor","Flags unsupported-generation risk. Lower is better."),
]

def add_slide_number(slide, num):
    box = slide.shapes.add_textbox(PptInches(12.3), PptInches(7.08), PptInches(0.7), PptInches(0.22))
    p = box.text_frame.paragraphs[0]; p.text = str(num); p.font.size = PptPt(8); p.font.color.rgb = PptRGBColor(110,110,110); p.alignment = PP_ALIGN.RIGHT
def add_metric_card(slide, x, y, w, h, title, value, subtitle=""):
    shape = slide.shapes.add_shape(1, PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = PptRGBColor(235,245,255); shape.line.color.rgb = PptRGBColor(160,190,220)
    tf=shape.text_frame; tf.clear()
    p=tf.paragraphs[0]; p.text=title; p.font.size=PptPt(10); p.font.bold=True; p.font.color.rgb=PptRGBColor(20,56,92)
    q=tf.add_paragraph(); q.text=value; q.font.size=PptPt(18); q.font.bold=True; q.font.color.rgb=PptRGBColor(0,97,120)
    if subtitle:
        r=tf.add_paragraph(); r.text=subtitle[:110]; r.font.size=PptPt(8); r.font.color.rgb=PptRGBColor(90,90,90)

def create_pptx(risk: pd.DataFrame, feeder: pd.DataFrame, impact: pd.DataFrame, rag_eval: pd.DataFrame) -> None:
    prs=Presentation(); prs.slide_width=PptInches(13.333); prs.slide_height=PptInches(7.5); blank=prs.slide_layouts[6]
    def slide(title, subtitle=None):
        s=prs.slides.add_slide(blank); bg=s.background.fill; bg.solid(); bg.fore_color.rgb=PptRGBColor(248,250,252); add_title(s,title,subtitle); add_slide_number(s,len(prs.slides)); return s
    vals=impact.set_index("metric")["value"]; annual=vals.get("Scaled annual recovery for 25,000-customer pilot", vals["Annual revenue recovery opportunity from high-risk customers"]); pilot_annual=vals["Annual revenue recovery opportunity from high-risk customers"]; fte=vals["Modeled maximum redeployable capacity"]; reduction=vals["Modeled repetitive workload reduction"]
    s=slide("Nepal Energy Agentic AI Command Center","Autonomous multi-agent demo for loss reduction, revenue recovery, workforce productivity, Agentic RAG, evaluation, and reporting")
    add_metric_card(s,0.7,1.3,2.8,1.15,"Annual recovery",money(annual),"Modeled 25,000-customer pilot"); add_metric_card(s,3.75,1.3,2.8,1.15,"Workload reduction",pct(reduction),"Repetitive manual work"); add_metric_card(s,6.8,1.3,2.8,1.15,"Redeployable capacity",f"{fte:.1f} FTE","Capacity, not layoffs"); add_metric_card(s,9.85,1.3,2.8,1.15,"Agents",f"{len(AGENT_BLUEPRINT)}","Coordinated workflow")
    add_bullets(s,0.9,3.05,11.7,3.3,["Best use case: NEA-style Energy Loss, Revenue Protection, and Workforce Productivity Command Center.","Not a chatbot: an autonomous workflow using data, agents, retrieval, evaluation, charts, PPTX, and DOCX.","All data is synthetic for safe public demonstration; production requires approved data and human approval gates."],17)
    s=slide("Why this use case wins for Nepal energy","Better than a chatbot because it targets measurable outcomes")
    add_bullets(s,0.7,1.15,12,5.7,["Revenue protection: flags suspicious customers and estimates recoverable energy/revenue.","Distribution operations: prioritizes feeders where losses, overload, suspicious accounts, and outage hours intersect.","Management productivity: turns manual spreadsheets and meeting preparation into automated evidence-based reporting.","Planning value: connects distribution loss control with hydropower, EV load, and power-trade context.","Stakeholder value: produces transparent RAG metrics, charts, PowerPoint, DOCX, and agent logs."],18)
    s=slide("Autonomous agent map","Agents communicate through shared data, evidence, metrics, charts, and reports")
    add_bullets(s,0.65,1.1,6,5.8,["Orchestrator controls the sequence and triggers each agent.","Data Agent prepares CSVs and validates tables.","Customer Risk Agent scores suspicious accounts.","Feeder Agent combines customer risk with feeder indicators.","Hydropower & Trade Agent adds seasonal and strategic planning context."],15)
    add_bullets(s,6.8,1.1,6,5.8,["Workforce Agent calculates manual-work reduction and redeployment capacity.","Agentic RAG Agent retrieves contextual evidence.","RAG Evaluation Agent checks quality and hallucination risk.","Visualization Agent generates charts.","Reporting Agent creates PPTX, DOCX, CSV, and JSON outputs."],15)
    s=slide("How agents reduce work by employee level","The goal is redeployment to higher-value work, not automatic replacement")
    add_bullets(s,0.7,1.1,12,5.9,["Executives: decision-ready summaries instead of repeated slide-update requests.","Managers: less manual feeder, field-visit, and monthly performance prioritization.","IT/MIS staff: less data cleaning, CSV merging, dashboard refresh, and report-generation work.","Revenue protection teams: focus on validated high-risk cases instead of scanning thousands of accounts.","Field supervisors: ranked inspection priorities and evidence packets while retaining final judgment.","Finance/audit teams: consistent revenue-recovery assumptions, traceable metrics, and logs."],16)
    s=slide("Business impact: revenue and workforce","Modeled demonstration numbers generated from synthetic data")
    add_bullets(s,0.7,1.15,5.8,5.5,[f"Pilot annual recovery: {money(pilot_annual)}",f"Scaled 25,000-customer annual recovery: {money(annual)}",f"Monthly recovery opportunity: {money(vals['Monthly revenue recovery opportunity'])}",f"Baseline repetitive work: {vals['Baseline repetitive manual work']:.1f} FTE-equivalent/month",f"Conservative half-workload target: {vals['Conservative half-workload redeployment target']:.1f} FTE-equivalent/month",f"Modeled max redeployable capacity: {vals['Modeled maximum redeployable capacity']:.1f} FTE-equivalent/month"],14)
    s.shapes.add_picture(str(CHART_REVENUE),PptInches(6.8),PptInches(1.2),width=PptInches(5.9))
    s=slide("Technology stack used in this demo","Simple public-cloud-friendly architecture; production upgrades are straightforward")
    for i,(layer,tool,why) in enumerate(TECH_STACK[:8]): add_metric_card(s,0.55+(i%2)*6.35,1.0+(i//2)*1.38,5.95,1.07,layer,tool,why)
    s=slide("RAG, LLM, embedding, and search design","Exact technology choices used in the demo")
    add_bullets(s,0.7,1.1,12,5.9,["LLM used: local deterministic demo LLM/template generator that writes answers only from retrieved evidence; no API key required.","RAG used: Agentic RAG because retrieval is one agent inside a larger autonomous workflow.","Hybrid search used: contextual lexical-semantic retrieval using TF-IDF n-grams plus topic/source/context prefixes.","Embedding used: contextual TF-IDF vector embedding for reliable Streamlit Community Cloud deployment.","Production embedding upgrade: Gemini text-embedding-004, BGE/E5, Instructor, or OpenAI embeddings.","Production vector DB upgrade: Chroma, FAISS, Pinecone, Weaviate, Milvus, or PostgreSQL pgvector.","Production LLM upgrade: Gemini, OpenAI, Groq, or local Ollama models with citation and grounding checks."],16)
    s=slide("Agentic RAG workflow","Question -> retrieval -> grounded answer -> evaluation -> reporting")
    add_bullets(s,0.7,1.1,12,5.9,["Knowledge CSV rows become contextual chunks with topic, source, URL, and Nepal energy context.","Chunks are embedded into a local vector matrix using TF-IDF n-grams over contextual text.","The RAG Agent retrieves top evidence chunks for each executive question.","The grounded generator builds answer text from retrieved evidence to reduce unsupported claims.","The Evaluation Agent calculates MRR, relevance, similarity, groundedness, and hallucination-risk estimate.","The Reporting Agent writes metrics and answer previews into PPTX, DOCX, CSV, and Streamlit views."],16)
    s=slide("RAG evaluation matrix","Quality metrics included in dashboard and reports")
    add_bullets(s,0.7,1.0,5.8,5.8,[f"{name}: {meaning}" for name,_,meaning in RAG_METRIC_DEFS],13); s.shapes.add_picture(str(CHART_RAG),PptInches(6.8),PptInches(1.2),width=PptInches(5.8))
    s=slide("RAG evaluation results","Average scores from synthetic demo questions")
    add_metric_card(s,0.75,1.35,2.2,1.25,"MRR",f"{rag_eval['mrr'].mean():.2f}","Higher is better"); add_metric_card(s,3.2,1.35,2.2,1.25,"Top-1 relevance",f"{rag_eval['relevance_top1'].mean():.2f}","Higher is better"); add_metric_card(s,5.65,1.35,2.2,1.25,"Similarity",f"{rag_eval['avg_similarity_top5'].mean():.2f}","Top-5 average"); add_metric_card(s,8.1,1.35,2.2,1.25,"Groundedness",f"{rag_eval['groundedness_estimate'].mean():.2f}","Higher is better"); add_metric_card(s,10.55,1.35,2.2,1.25,"Hallucination risk",f"{rag_eval['hallucination_rate_estimate'].mean():.2f}","Lower is better")
    add_bullets(s,0.8,3.2,11.8,3.0,["These are demonstration metrics from synthetic evidence and deterministic retrieval.","For production, the same evaluation layer should run against human-labeled questions and field outcomes.","Grounding and hallucination monitoring help managers decide where human review is mandatory."],16)
    s=slide("Suspicious activity detection","Customer-level anomaly scoring and revenue recovery"); s.shapes.add_picture(str(CHART_RISK),PptInches(0.7),PptInches(1.15),width=PptInches(5.9)); s.shapes.add_picture(str(CHART_REVENUE),PptInches(6.85),PptInches(1.15),width=PptInches(5.8))
    s=slide("Feeder prioritization","Combines feeder losses, overload, outage hours, and suspicious customer density"); s.shapes.add_picture(str(CHART_FEEDER),PptInches(0.75),PptInches(1.1),width=PptInches(6.1)); add_bullets(s,7.1,1.3,5.6,4.9,["Managers see which feeders deserve inspection first.","Revenue teams can coordinate field visits with feeder-loss priorities.","IT teams automate monthly ranking instead of ad-hoc spreadsheets.","Executive dashboards track progress by feeder, district, and center."],16)
    s=slide("Workforce productivity and redeployment","Automation reduces repetitive work while increasing governance and field effectiveness"); s.shapes.add_picture(str(CHART_WORKFORCE),PptInches(0.8),PptInches(1.15),width=PptInches(6)); add_bullets(s,7,1.25,5.7,5,[f"Modeled repetitive workload reduction: {pct(reduction)}.",f"Redeployable capacity estimate: {fte:.1f} FTE-equivalent/month.","Best framing: reduce manual reporting, repeated data cleaning, and low-value screening.","Employees shift toward validation, exception handling, customer engagement, safety, and governance."],16)
    s=slide("PowerPoint, DOCX, and Streamlit automation","A single run produces stakeholder-ready outputs"); add_bullets(s,0.7,1.1,12,5.9,["Command: python main.py generates synthetic CSVs, agent scores, RAG metrics, charts, PPTX, DOCX, logs, and running steps.","Command: streamlit run streamlit_app.py launches the public dashboard.","Streamlit Community Cloud main file: streamlit_app.py.","The same generated CSVs/images are reused in dashboard, Word report, and presentation for consistency.","The flat project is easy to teach, modify, upload to GitHub, and deploy publicly."],17)
    s=slide("Production roadmap and governance","How this demo becomes a real NEA pilot"); add_bullets(s,0.7,1.1,12,5.9,["Phase 1: offline pilot using approved billing, meter, feeder, outage, and complaint data.","Phase 2: validate risk scores with field inspection outcomes and recovered revenue.","Phase 3: role-based dashboards for executives, managers, IT, revenue protection, and field supervisors.","Phase 4: audit trails, cybersecurity, privacy, explainability, model monitoring, and human approval gates.","Phase 5: scale after transparent evaluation and stakeholder training."],17)
    s=slide("Decision ask for stakeholders","Start with a focused, measurable, low-risk pilot"); add_bullets(s,0.7,1.1,12,5.9,["Approve a 60-90 day pilot in selected distribution centers using historical data.","Measure recovered revenue, inspection hit rate, repetitive reporting hours saved, and RAG answer quality.","Use AI as decision support and productivity support, not as an unsupervised decision maker.","Keep humans responsible for approvals, customer actions, field safety, and legal decisions."],18)
    s=slide("Important disclaimer","Synthetic demo for learning and stakeholder discussion"); add_bullets(s,0.7,1.2,12,5.8,["All datasets and results are synthetic and created for demonstration.","Revenue and employee-capacity numbers are modeled estimates, not official NEA forecasts.","Employee impact is framed as workload reduction and capacity redeployment, not automatic layoffs.","Real deployment requires approved data access, privacy review, cybersecurity review, legal review, and human controls."],18)
    prs.save(PPTX_FILE)

def set_doc_style(doc: Document):
    styles=doc.styles; styles["Normal"].font.name="Aptos"; styles["Normal"].font.size=Pt(10.5)
    for style_name in ["Heading 1","Heading 2","Heading 3"]:
        style=styles[style_name]; style.font.name="Aptos Display"; style.font.color.rgb=RGBColor(20,56,92)
def add_cover(doc: Document):
    title=doc.add_paragraph(); title.alignment=WD_ALIGN_PARAGRAPH.CENTER
    run=title.add_run("Nepal Energy Agentic AI Command Center"); run.bold=True; run.font.size=Pt(24); run.font.color.rgb=RGBColor(20,56,92)
    sub=doc.add_paragraph("Autonomous multi-agent demo for loss reduction, revenue recovery, workforce productivity, Agentic RAG, evaluation, and stakeholder reporting"); sub.alignment=WD_ALIGN_PARAGRAPH.CENTER
    doc.add_paragraph(f"Generated: {TODAY}").alignment=WD_ALIGN_PARAGRAPH.CENTER
    doc.add_paragraph("This report is designed for NEA-style stakeholders, executives, managers, IT teams, revenue protection teams, and public-sector AI decision makers. All data is synthetic and safe for demonstration.")
    doc.add_page_break()
def add_section_page(doc: Document, title: str, paragraphs: List[str], bullets: List[str] | None=None, chart: Path | None=None, page_break: bool=True):
    doc.add_heading(title, level=1)
    for para in paragraphs: doc.add_paragraph(para)
    if bullets:
        for b in bullets: doc.add_paragraph(b, style="List Bullet")
    if chart and chart.exists(): doc.add_picture(str(chart), width=Inches(6.25))
    if page_break: doc.add_page_break()
def add_df_table(doc: Document, df: pd.DataFrame, title: str, max_rows: int=10):
    doc.add_heading(title, level=2); show=df.head(max_rows).copy()
    table=doc.add_table(rows=1, cols=len(show.columns)); table.style="Light Grid Accent 1"
    for j,c in enumerate(show.columns): table.rows[0].cells[j].text=str(c)
    for _,row in show.iterrows():
        cells=table.add_row().cells
        for j,c in enumerate(show.columns): cells[j].text=str(row[c])[:120]
def create_docx(risk: pd.DataFrame, feeder: pd.DataFrame, impact: pd.DataFrame, rag_eval: pd.DataFrame) -> None:
    doc=Document(); set_doc_style(doc)
    sec=doc.sections[0]; sec.top_margin=Inches(0.65); sec.bottom_margin=Inches(0.65); sec.left_margin=Inches(0.7); sec.right_margin=Inches(0.7)
    add_cover(doc)
    vals=impact.set_index("metric")["value"]; annual=vals.get("Scaled annual recovery for 25,000-customer pilot", vals["Annual revenue recovery opportunity from high-risk customers"])

    # 1
    add_section_page(doc,"1. Executive summary",[
        "The recommended use case is an NEA-style Energy Loss, Revenue Protection, and Workforce Productivity Command Center. It is stronger than a chatbot because it performs a full autonomous workflow: data preparation, suspicious activity scoring, feeder ranking, contextual RAG evidence retrieval, quality evaluation, charting, and stakeholder reporting.",
        f"The synthetic demo estimates a scaled annual revenue recovery opportunity of {money(annual)} and a repetitive workload reduction of {pct(vals['Modeled repetitive workload reduction'])}. These numbers are modeled demonstration outputs, not official forecasts.",
        "Employee impact is framed as staff-capacity redeployment rather than direct layoffs. Managers, IT staff, revenue-protection analysts, and field supervisors remain essential for approvals, governance, inspections, safety, and customer communication."
    ])

    # 2
    doc.add_heading("2. Business impact dashboard",level=1)
    table=doc.add_table(rows=1, cols=4); table.style="Light Shading Accent 1"
    for i,h in enumerate(["Metric","Value","Unit","Explanation"]): table.rows[0].cells[i].text=h
    for _,r in impact.iterrows():
        row=table.add_row().cells; row[0].text=str(r["metric"]); value=float(r["value"])
        row[1].text=f"{value:,.1f}" if abs(value)<100000 else f"{value:,.0f}"; row[2].text=str(r["unit"]); row[3].text=str(r["explanation"])[:140]
    doc.add_paragraph("This table gives stakeholders concrete synthetic pilot-style assumptions for financial impact, operational impact, and workforce productivity.")
    doc.add_page_break()

    # 3
    doc.add_heading("3. How many agents are used?",level=1)
    doc.add_paragraph(f"The demo uses {len(AGENT_BLUEPRINT)} agents. They are not isolated chatbots. They communicate through a shared workflow controlled by the Orchestrator Agent.")
    for num,name,role,work,employee,reduce,comm in AGENT_BLUEPRINT:
        doc.add_heading(f"{num}. {name}", level=2)
        for label, text in [("Role",role),("Main work",work),("Employee level supported",employee),("Work reduced",reduce),("Communication",comm)]:
            doc.add_paragraph(f"{label}: {text}", style="List Bullet")
    doc.add_page_break()

    pages = [
        ("4. Agent communication model", ["Agents communicate sequentially and through shared artifacts. The Orchestrator starts the run, the Data Agent prepares CSVs, the Risk Agent scores customers, the Feeder Agent combines results, the RAG Agent retrieves evidence, the Evaluation Agent measures quality, and the Reporting Agent creates final outputs.", "Some agents can work independently, but the highest value comes from communication. Feeder ranking improves after receiving suspicious-customer density from the Risk Agent."], ["Independent work: data validation, RAG indexing, chart generation, workforce scoring.","Collaborative work: feeder prioritization, revenue recovery, RAG evaluation, and report generation.","Production upgrade: LangGraph, CrewAI, AutoGen, Google ADK, or custom orchestration."], None),
        ("5. Workforce impact by employee level", ["The demo estimates repetitive workload reduction in FTE-equivalent terms. This helps show how low-value repetitive work can be automated so employees can be redeployed to higher-value tasks.", f"Modeled maximum redeployable capacity: {vals['Modeled maximum redeployable capacity']:.1f} FTE-equivalent/month. Conservative half-workload redeployment target: {vals['Conservative half-workload redeployment target']:.1f} FTE-equivalent/month."], ["Managerial level: less manual performance compilation and feeder prioritization.","IT/MIS level: less CSV cleaning, dashboard refresh, and ad-hoc data pulling.","Revenue-protection level: less manual account screening.","Field-supervisor level: better inspection priorities and evidence packets.","Finance/audit level: repeatable assumptions and logs."], CHART_WORKFORCE),
        ("6. Revenue generation and recovery logic", ["Revenue generation comes from prioritizing high-risk accounts and high-loss feeders. The Customer Risk Agent estimates recovery opportunity, then the Business Impact Agent scales the pilot.", "In production, revenue numbers should be validated against field inspections, corrected meter data, recovered arrears, recovered kWh, and reduced losses."], [f"Modeled annual recovery opportunity: {money(annual)}.",f"Modeled monthly recovery opportunity: {money(vals['Monthly revenue recovery opportunity'])}.","Validation KPI: inspection hit rate, recovered kWh, recovered NPR, reduced feeder loss, reduced repeat complaints."], CHART_REVENUE),
        ("7. Technology architecture overview", ["The project is intentionally simple for public demonstration and Streamlit Community deployment. It uses a flat Python structure and does not require paid API keys.", "The architecture still demonstrates multi-agent orchestration, Agentic RAG, contextual embeddings, hybrid retrieval, model evaluation, and automated stakeholder reporting."], None, None),
    ]
    for title, paras, bullets, chart in pages:
        add_section_page(doc,title,paras,bullets,chart)

    # 8 tech stack as bullets, not a huge table
    doc.add_heading("8. Detailed technology stack", level=1)
    for layer,tool,why in TECH_STACK:
        doc.add_heading(layer, level=2)
        doc.add_paragraph(f"Technology used: {tool}", style="List Bullet")
        doc.add_paragraph(f"Why it is used: {why}", style="List Bullet")
    doc.add_page_break()

    more_pages = [
        ("9. Which LLM is used?", ["This public demo uses a local deterministic demo LLM pattern instead of a paid hosted LLM. The answer generator creates stakeholder-friendly grounded answers directly from retrieved evidence chunks.", "For production, the same interface can be upgraded to Gemini, OpenAI, Groq, or Ollama models."], ["Demo mode: deterministic grounded generator from retrieved chunks.","Low-cost option: Groq or Gemini free/low-cost tier.","Local option: Ollama with Llama or Mistral models.","Enterprise option: Vertex AI Gemini, Azure OpenAI, OpenAI API, or private endpoint with audit logging."], None),
        ("10. Which RAG is used?", ["The demo uses Agentic RAG. A simple chatbot retrieves text and answers a question. Agentic RAG retrieves evidence as part of a broader autonomous workflow that includes scoring, decision support, quality evaluation, and reporting.", "The RAG Agent is one member of the agent team, and its retrieved evidence supports executive answers, business justification, and transparent evaluation metrics."], ["RAG input: synthetic Nepal energy knowledge CSV.","Chunking: knowledge rows are split into small context-style chunks.","Context enrichment: every chunk is prefixed with topic, source, and Nepal energy context.","Retrieval: top chunks are selected using local vector similarity.","Generation: answers are formed from retrieved evidence."], None),
        ("11. Which embedding and hybrid search are used?", ["The demo uses contextual TF-IDF vector embeddings. This lightweight local embedding approach is easy to run and stable for public deployment.", "The hybrid-search idea combines lexical n-gram similarity with contextual prefixes such as topic, source, and Nepal energy context."], ["Current embedding: contextual TF-IDF n-gram vectors.","Current hybrid retrieval: keyword/phrase matching plus context-enriched similarity.","Production embedding options: text-embedding-004, BGE, E5, Instructor, OpenAI embeddings.","Production hybrid search options: BM25 + dense embeddings + reranking.","Production vector databases: Chroma, FAISS, Pinecone, Weaviate, Milvus, pgvector."], None),
        ("12. RAG evaluation metrics", ["RAG evaluation helps stakeholders understand whether AI is retrieving the right evidence and whether answers are grounded. This demo calculates MRR, relevance, similarity, groundedness, and hallucination-risk estimates."], [f"{name} ({full}): {meaning}" for name,full,meaning in RAG_METRIC_DEFS], CHART_RAG),
    ]
    for title, paras, bullets, chart in more_pages:
        add_section_page(doc,title,paras,bullets,chart)

    # 13 RAG table
    doc.add_heading("13. RAG evaluation results table", level=1)
    t=doc.add_table(rows=1, cols=6); t.style="Light Grid Accent 1"
    for i,h in enumerate(["Question","MRR","Relevance","Similarity","Groundedness","Hallucination risk"]): t.rows[0].cells[i].text=h
    for _,r in rag_eval.iterrows():
        cells=t.add_row().cells
        vals2=[r["question"],r["mrr"],r["relevance_top1"],r["avg_similarity_top5"],r["groundedness_estimate"],r["hallucination_rate_estimate"]]
        for i,v in enumerate(vals2): cells[i].text=str(v)
    doc.add_paragraph("These metrics are demonstration estimates. In a real deployment, questions should be created with NEA experts and reviewed monthly.")
    doc.add_page_break()

    # 14-23
    tail_pages = [
        ("14. Suspicious customer activity detection", ["The Customer Risk Agent combines unsupervised anomaly detection and interpretable business rules. Isolation Forest identifies unusual consumption patterns while rules add practical signals such as meter tamper events, arrears, complaint count, and consumption drops.", "The goal is not just to identify outliers; the goal is to prioritize cases useful for revenue-protection teams."], ["Input: synthetic customer consumption, billing, meter, arrears, and complaint data.","Output: risk score, action recommendation, and estimated annual recovery.","Employee impact: revenue teams review fewer low-value cases."], CHART_RISK),
        ("15. Top suspicious customer examples", ["The table below shows examples of the highest-scored synthetic customer records. These are not real customers. They demonstrate how the report can provide a transparent inspection shortlist."], None, None),
        ("16. Feeder prioritization", ["The Feeder Prioritization Agent identifies which feeders deserve attention first. It combines loss percentage, outage hours, overload percentage, suspicious customer count, and recovery opportunity.", "This helps managers allocate limited field resources more effectively."], ["Manager impact: better crew planning.","IT impact: fewer monthly spreadsheet ranking requests.","Field impact: clearer inspection priorities and evidence packages."], CHART_FEEDER),
        ("17. Top feeder priority examples", ["The table below shows the highest-priority synthetic feeders. In a real deployment, these outputs should be validated against field inspections and measured recovery outcomes."], None, None),
        ("18. PowerPoint and DOCX automation", ["The Reporting Agent automatically generates both the PowerPoint and DOCX report. This reduces the repeated work of slide preparation, document formatting, and monthly management reporting.", "The same results power the dashboard, charts, CSVs, PPTX, and DOCX."], ["Output: stakeholder PowerPoint.","Output: long stakeholder DOCX report.","Output: risk CSV, feeder CSV, impact CSV, RAG CSV, agent log JSON.","Output: charts reused across reports and dashboard."], None),
        ("19. Streamlit Community Cloud deployment", ["The project is designed for Streamlit Community Cloud. The repository should contain main.py, streamlit_app.py, requirements.txt, CSV files, generated charts, and reports directly in the root.", "Because the demo requires no paid API key, people in Nepal can open the public URL and use the dashboard without secrets."], ["Local run: python main.py, then streamlit run streamlit_app.py.","Cloud run: upload flat folder to GitHub and set streamlit_app.py as entry point.","Public demo mode: safe synthetic data only.","Production mode: approved data, authentication, logging, human approval gates."], None),
        ("20. Production upgrade path", ["The simple demo can become an enterprise pilot by replacing synthetic CSVs with approved NEA data, replacing TF-IDF with dense embeddings, adding a vector DB, adding a production LLM, and adding governance controls.", "The most important production upgrade is governance: explainability, role-based access, audit logs, privacy, and field-outcome validation."], ["Data: billing, meter events, SCADA/feeder data, outage records, complaints, inspections.","Retrieval: BM25 + dense embeddings + reranker.","Vector store: Chroma/FAISS for pilot; pgvector/Pinecone/Weaviate/Milvus for production.","Agents: LangGraph, CrewAI, AutoGen, Google ADK, or custom orchestration.","Monitoring: RAG metrics, model drift, human feedback, recovered revenue, inspection hit rate."], None),
        ("21. Governance, safety, and human control", ["Energy-sector AI should not make punitive customer decisions automatically. The system should recommend priorities, explain reasons, and provide evidence.", "The demo treats AI as decision support to improve focus, reduce repetitive work, and increase transparency."], ["Use role-based access for executives, managers, IT, field teams, and auditors.","Keep audit logs for every score, retrieved evidence, and report version.","Use explainable risk factors rather than black-box decisions alone.","Protect privacy and separate demo data from real customer data."], None),
        ("22. Public context and references", ["The demo is grounded in public energy-sector context and open-source feasibility patterns. Links are stakeholder references and are not called automatically by the app."], [f"{ref['name']}: {ref['note']} URL: {ref['url']}" for ref in PUBLIC_CONTEXT], None),
        ("23. Recommended pilot plan", ["A practical NEA-style pilot should start small, focus on measurable outcomes, and avoid overpromising. A 60-90 day pilot can test inspection hit rate, recovered revenue, reporting speed, and management visibility.", "The pilot should include baseline metrics before AI, weekly review, field validation, staff feedback, and monthly leadership reporting."], ["Pilot area: selected distribution centers or feeders.","Pilot data: historical billing, meter events, feeder losses, outages, complaints, inspections.","Pilot KPIs: recovered NPR, recovered kWh, inspection hit rate, workload hours saved, RAG groundedness, hallucination risk.","Pilot governance: human approval, legal review, cybersecurity review, customer communication policy."], None),
        ("24. Final stakeholder message", ["The strongest stakeholder message is that this is an AI operating workflow, not a chatbot. It turns data into prioritized action, documents into grounded evidence, manual reporting into automated deliverables, and AI uncertainty into evaluation metrics.", "The safest way to present employee impact is to say that the system can reduce repetitive manual workload by roughly half or more in the modeled workflow, allowing redeployment to supervision, validation, field action, customer service, and governance."], ["Start with a controlled pilot.","Measure real business outcomes.","Keep humans in charge.","Scale only after transparent validation."], None),
    ]
    for title, paras, bullets, chart in tail_pages:
        add_section_page(doc,title,paras,bullets,chart,page_break=not title.startswith("24."))
        if title.startswith("15."):
            simple=risk[["customer_id","district","tariff_category","feeder_id","risk_score_0_100","estimated_revenue_recovery_npr_year","recommended_action"]].head(8)
            add_df_table(doc,simple,"Top risk-scored customers",8); doc.add_page_break()
        if title.startswith("17."):
            simplef=feeder[["feeder_id","district","priority_score_0_100","technical_nontechnical_loss_pct","peak_loading_ratio","suspicious_customers","revenue_opportunity_npr","field_visit_recommendation"]].head(8)
            add_df_table(doc,simplef,"Top feeder priorities",8); doc.add_page_break()
    doc.save(DOCX_FILE)

def create_running_steps() -> None:
    README_OUTPUT.write_text(r"""# Nepal Energy Agentic AI Demo - Simple Flat Version

This version is intentionally simple: no nested code folders, no package imports, and no secrets required.

## Files

- `main.py` - Generates data, runs agents, RAG, evaluation, charts, PPTX, DOCX, and CSV outputs.
- `streamlit_app.py` - Streamlit dashboard.
- `requirements.txt` - Python packages.
- `synthetic_*.csv` - Synthetic Nepal energy data.
- Generated files: `nepal_energy_agentic_ai_demo.pptx`, `nepal_energy_agentic_ai_report.docx`, `risk_scores.csv`, `rag_evaluation.csv`, charts, and logs.

## Local run on Windows PowerShell

```powershell
cd C:\Users\Parshuram\Downloads\nepal_energy_flat_demo
python -m venv .venv
Set-ExecutionPolicy -Scope Process -ExecutionPolicy Bypass
.\.venv\Scripts\Activate.ps1
pip install -r requirements.txt
python main.py
streamlit run streamlit_app.py
```

If activation gives trouble, skip activation and run:

```powershell
.\.venv\Scripts\python.exe -m pip install -r requirements.txt
.\.venv\Scripts\python.exe main.py
.\.venv\Scripts\python.exe -m streamlit run streamlit_app.py
```

## Streamlit Community Cloud deployment

1. Create a GitHub repository.
2. Upload all files from this folder to the repository root.
3. Go to Streamlit Community Cloud.
4. Select your GitHub repository.
5. Set the main file path to: `streamlit_app.py`.
6. Deploy.

No API key is required for demo mode.

## Important note

All data is synthetic. Employee impact is shown as repetitive workload reduction and capacity redeployment, not direct layoffs.
""", encoding="utf-8")


def run_autonomous_pipeline() -> Dict:
    logs: List[Dict] = []
    log_event(logs, "Orchestrator Agent", "Starting full autonomous Nepal energy AI pipeline.")

    create_synthetic_csvs()
    log_event(logs, "Data Agent", "Created/updated synthetic CSV files.", {"files": [p.name for p in [CUSTOMERS_CSV, FEEDERS_CSV, HYDRO_CSV, WORKFORCE_CSV, KNOWLEDGE_CSV]]})

    customers = pd.read_csv(CUSTOMERS_CSV)
    feeders = pd.read_csv(FEEDERS_CSV)
    hydro = pd.read_csv(HYDRO_CSV)
    workforce = pd.read_csv(WORKFORCE_CSV)
    knowledge = pd.read_csv(KNOWLEDGE_CSV)

    risk = run_customer_risk(customers)
    risk.to_csv(RISK_CSV, index=False)
    log_event(logs, "Risk Detection Agent", "Scored suspicious customer activity and revenue recovery opportunity.", {"high_risk_customers": int((risk["risk_score_0_100"] >= 70).sum())})

    feeder_priority = run_feeder_priority(feeders, risk)
    feeder_priority.to_csv(FEEDER_PRIORITY_CSV, index=False)
    log_event(logs, "Feeder Agent", "Prioritized feeders for inspection and digital monitoring.", {"top_feeder": feeder_priority.iloc[0]["feeder_id"]})

    index = build_rag_index(knowledge)
    rag_eval = evaluate_rag(index)
    rag_eval.to_csv(RAG_EVAL_CSV, index=False)
    log_event(logs, "Agentic RAG Agent", "Built contextual retrieval index and generated RAG evaluation.", {"chunks": int(len(index.chunks)), "avg_mrr": float(rag_eval["mrr"].mean())})

    impact = calculate_business_impact(risk, workforce, hydro)
    impact.to_csv(BUSINESS_IMPACT_CSV, index=False)
    log_event(logs, "Business Impact Agent", "Calculated revenue recovery and workforce redeployment estimates.", {"annual_revenue_npr": float(impact.iloc[0]["value"])})

    make_charts(risk, feeder_priority, impact, rag_eval)
    log_event(logs, "Visualization Agent", "Generated charts.")

    create_pptx(risk, feeder_priority, impact, rag_eval)
    create_docx(risk, feeder_priority, impact, rag_eval)
    create_running_steps()
    log_event(logs, "Reporting Agent", "Created stakeholder PPTX, DOCX, CSV outputs, and running steps.")

    AGENT_LOG_JSON.write_text(json.dumps(logs, indent=2), encoding="utf-8")
    log_event(logs, "Orchestrator Agent", "Completed full autonomous run.")

    return {
        "risk_rows": len(risk),
        "top_customer": risk.iloc[0]["customer_id"],
        "annual_revenue_recovery_npr": float(impact.iloc[0]["value"]),
        "workload_reduction_pct": float(impact.query("metric == 'Modeled repetitive workload reduction'")["value"].iloc[0]),
        "redeployable_fte": float(impact.query("metric == 'Modeled maximum redeployable capacity'")["value"].iloc[0]),
        "pptx": str(PPTX_FILE),
        "docx": str(DOCX_FILE),
    }


if __name__ == "__main__":
    summary = run_autonomous_pipeline()
    print("\nSUCCESS: Nepal Energy Agentic AI demo completed.\n")
    for k, v in summary.items():
        if isinstance(v, float):
            print(f"{k}: {v:,.2f}")
        else:
            print(f"{k}: {v}")
